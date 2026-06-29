#!/usr/bin/env python3
"""
Convert a folder of slide-N.png images into a PPTX file.
Each PNG becomes a full-bleed image on a blank slide.
Usage: python3 pngs_to_pptx.py <png_dir> <output.pptx>
"""
import sys, os, glob
from pptx import Presentation
from pptx.util import Inches, Emu

SW = 10.0
SH = 5.625

def build(png_dir, out_path):
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    pngs = sorted(glob.glob(os.path.join(png_dir, 'slide-*.png')),
                  key=lambda p: int(os.path.basename(p).replace('slide-','').replace('.png','')))

    if not pngs:
        print(f'[pptx] No slide-*.png found in {png_dir}', file=sys.stderr)
        sys.exit(1)

    for i, png in enumerate(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
        print(f'[pptx] slide {i+1}/{len(pngs)}: {os.path.basename(png)}', file=sys.stderr)

    prs.save(out_path)
    print(f'[pptx] Saved {len(pngs)} slides → {out_path}')

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: pngs_to_pptx.py <png_dir> <output.pptx>', file=sys.stderr)
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])
