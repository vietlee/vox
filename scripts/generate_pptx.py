#!/usr/bin/env python3
"""
VOX Slide PPTX Generator — element-based deck schema.
Reads deck JSON (compiled by ContentOutlineGenerator#build_deck_schema).
All element coords in inches; uses python-pptx directly.
"""
import sys, json, math, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I

try:
    from pptx.oxml.ns import qn
    import lxml.etree as etree
except ImportError:
    qn = None

SW = 10.0   # slide width in inches  (same as Ruby compiler)
SH = 5.625  # slide height in inches

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def safe_color(s): return hex_to_rgb(s) if s and s.startswith('#') else RGBColor(0x1F,0x2A,0x44)

def add_shape_fill(shape, color_hex, opacity=1.0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)

def add_text_box(slide, x, y, w, h, text, style):
    from pptx.util import Pt, Emu
    txBox = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    align = style.get('align','left')
    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.RIGHT if align == 'right' else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text or ''
    run.font.size = Pt(style.get('fontSize', 11))
    run.font.bold = (style.get('fontWeight', 400) >= 700)
    run.font.italic = style.get('fontStyle','') == 'italic'
    col = style.get('color','#1F2A44')
    run.font.color.rgb = safe_color(col)
    run.font.name = 'Nunito' if style.get('fontFamily','') == 'Nunito' else 'Calibri'

def add_rounded_rect(slide, x, y, w, h, fill_hex, radius_pct=0.05, opacity=1.0, stroke_hex=None, stroke_w=0):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        I(x), I(y), I(w), I(h)
    )
    add_shape_fill(shape, fill_hex, opacity)
    if stroke_hex:
        shape.line.color.rgb = safe_color(stroke_hex)
        shape.line.width = Pt(stroke_w * 72) if stroke_w else Pt(1)
    else:
        shape.line.fill.background()
    # Apply rounding via XML
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is None:
            prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for av in avLst.findall(qn('a:gd')):
            avLst.remove(av)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name','adj')
        r = int(min(radius_pct, 0.5) * 100000) if isinstance(radius_pct, float) else 8000
        gd.set('fmla', f'val {r}')
    except Exception:
        pass
    return shape

def add_ellipse(slide, x, y, w, h, fill_hex, opacity=1.0):
    shape = slide.shapes.add_shape(9, I(x), I(y), I(w), I(h))  # 9 = oval
    add_shape_fill(shape, fill_hex, opacity)
    shape.line.fill.background()
    if opacity < 1.0:
        try:
            sp = shape._element
            spPr = sp.find(qn('p:spPr'))
            solidFill = spPr.find('.//' + qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = etree.SubElement(srgb, qn('a:alpha'))
                    alpha.set('val', str(int(opacity * 100000)))
        except Exception:
            pass
    return shape

ICON_CHARS = {
    'star': '★', 'check': '✓', 'rocket': '🚀', 'chart': '📊', 'money': '💰',
    'person': '👤', 'people': '👥', 'leaf': '🌿', 'lightbulb': '💡',
    'shield': '🛡', 'target': '🎯', 'crown': '♛', 'handshake': '🤝',
    'globe': '🌐', 'megaphone': '📣', 'search': '🔍', 'code': '⌨',
    'clock': '⏱', 'store': '🏪', 'percent': '%', 'heart': '❤',
    'truck': '🚚', 'phone': '📞',
}

def add_icon_box(slide, x, y, size, fill_hex, icon_name='star', icon_color='#FFFFFF'):
    shape = slide.shapes.add_shape(9, I(x), I(y), I(size), I(size))
    add_shape_fill(shape, fill_hex)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    ch = ICON_CHARS.get(icon_name, '●')
    run.text = ch
    run.font.size = Pt(size * 72 * 0.45)  # 45% of shape size in pts
    try:
        run.font.color.rgb = safe_color(icon_color)
    except Exception:
        pass
    return shape

def render_element(slide, el, slide_idx):
    el_type = el.get('type','')
    x = el.get('x', 0); y = el.get('y', 0)
    w = el.get('w', 1); h = el.get('h', 0.5)
    s = el.get('style', {})

    # Clip to slide bounds for visible elements
    if x > SW or y > SH or x + w < -1: return
    if el_type == 'text':
        add_text_box(slide, x, y, w, h, el.get('content',''), s)

    elif el_type in ('rect',):
        fill = s.get('fill','#E8F0FF')
        r = s.get('borderRadius', 0)
        radius_pct = 0.08 if r and r != 0 and r != '0' else 0
        stroke = s.get('stroke')
        sw2 = s.get('strokeWidth', 0)
        op = s.get('opacity', 1.0)
        add_rounded_rect(slide, x, y, w, h, fill, radius_pct, op, stroke, sw2)

    elif el_type == 'ellipse':
        fill = s.get('fill','#6366f1')
        op = s.get('opacity', 1.0)
        add_ellipse(slide, x, y, w, h, fill, op)

    elif el_type == 'line':
        from pptx.util import Pt
        line_shape = slide.shapes.add_shape(9, I(x), I(y), I(w), I(max(h, 0.03)))
        add_shape_fill(line_shape, s.get('stroke','#6366f1'))
        line_shape.line.fill.background()

    elif el_type == 'icon':
        bg = s.get('bgColor', s.get('fill','#6366f1'))
        ic = s.get('color', '#FFFFFF')
        add_icon_box(slide, x, y, w, bg, el.get('icon','star'), ic)

    elif el_type == 'chart_bar':
        render_bar_chart(slide, el)

    elif el_type == 'chart_donut':
        render_donut_chart(slide, el)

def set_chart_datalabels(chart_elem, show_val=True, show_cat=False, show_pct=False, font_size_pt=9):
    """Set data labels at chart level via XML — overrides series-level defaults."""
    try:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        # Work on the chart XML element (c:chart)
        plotArea = chart_elem.find('.//' + qn('c:plotArea'))
        if plotArea is None:
            return
        # Remove existing chart-level dLbls
        for old in plotArea.findall(qn('c:dLbls')):
            plotArea.remove(old)
        # Add fresh dLbls
        dLbls = etree.SubElement(plotArea, qn('c:dLbls'))
        numFmt = etree.SubElement(dLbls, qn('c:numFmt'))
        numFmt.set('formatCode', 'General'); numFmt.set('sourceLinked', '0')
        # Text style: font size
        txPr = etree.SubElement(dLbls, qn('c:txPr'))
        bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
        etree.SubElement(txPr, qn('a:lstStyle'))
        p = etree.SubElement(txPr, qn('a:p'))
        pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = etree.SubElement(pPr, qn('a:defRPr'))
        defRPr.set('sz', str(int(font_size_pt * 100)))  # hundredths of a point
        # Show flags — ALL must be explicit to avoid Office defaults
        for tag, val in [
            ('c:showLegendKey', '0'), ('c:showVal', '1' if show_val else '0'),
            ('c:showCatName', '1' if show_cat else '0'), ('c:showSerName', '0'),
            ('c:showPercent', '1' if show_pct else '0'), ('c:showBubbleSize', '0'),
        ]:
            e = etree.SubElement(dLbls, qn(tag)); e.set('val', val)
    except Exception as ex:
        print(f'[PPTX] datalabels error: {ex}', file=sys.stderr)

def set_series_color(series, hex_color):
    """Set solid fill on all data points and the series itself."""
    try:
        rgb = safe_color(hex_color)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = rgb
    except Exception as ex:
        print(f'[PPTX] series color error: {ex}', file=sys.stderr)

def set_point_colors(series, colors):
    """Color individual data points (bars) with per-point fills via XML."""
    try:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        ser_el = series._element
        for idx, hex_c in enumerate(colors):
            # Find or create dPt for this index
            dPt = None
            for existing in ser_el.findall(qn('c:dPt')):
                idx_el = existing.find(qn('c:idx'))
                if idx_el is not None and idx_el.get('val') == str(idx):
                    dPt = existing; break
            if dPt is None:
                dPt = etree.SubElement(ser_el, qn('c:dPt'))
                idx_el = etree.SubElement(dPt, qn('c:idx'))
                idx_el.set('val', str(idx))
            # Set fill
            spPr = dPt.find(qn('c:spPr'))
            if spPr is None:
                spPr = etree.SubElement(dPt, qn('c:spPr'))
            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is not None:
                spPr.remove(solidFill)
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgb.set('val', hex_c.lstrip('#'))
    except Exception as ex:
        print(f'[PPTX] point colors error: {ex}', file=sys.stderr)

def render_bar_chart(slide, el):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    x = el.get('x',0); y = el.get('y',0); w = el.get('w',4); h = el.get('h',2)
    chart_info = el.get('chart', {})
    data_arr = chart_info.get('data', [])
    if not data_arr: return
    # Get theme colors from the deck global (passed via el's theme_colors field or fallback)
    theme_colors = el.get('_theme_colors') or ['#6366F1','#8B5CF6','#A78BFA','#4F46E5','#7C3AED','#6D28D9']

    chart_data = ChartData()
    chart_data.categories = [str(d.get('label','')) for d in data_arr]
    chart_data.add_series('', [float(d.get('value', 0)) for d in data_arr])

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, I(x), I(y), I(w), I(h), chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    # Remove chart border
    try:
        graphic_frame.line.width = 0
    except Exception:
        pass

    series = chart.series[0]
    # Base series color = first theme color
    set_series_color(series, theme_colors[0])
    # Per-point colors (cycling through theme palette)
    point_colors = [theme_colors[i % len(theme_colors)] for i in range(len(data_arr))]
    set_point_colors(series, point_colors)

    # Data labels — value only, no category
    set_chart_datalabels(chart._element, show_val=True, show_cat=False)

def render_donut_chart(slide, el):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    x = el.get('x',0); y = el.get('y',0); w = el.get('w',4); h = el.get('h',2)
    chart_info = el.get('chart', {})
    data_arr = chart_info.get('data', [])
    if not data_arr: return
    theme_colors = el.get('_theme_colors') or ['#6366F1','#8B5CF6','#A78BFA','#4F46E5','#7C3AED','#6D28D9']

    chart_data = ChartData()
    chart_data.categories = [str(d.get('label','')) for d in data_arr]
    chart_data.add_series('', [float(d.get('value', 0)) for d in data_arr])

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, I(x), I(y), I(w), I(h), chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.has_title = False

    series = chart.series[0]
    set_series_color(series, theme_colors[0])
    point_colors = [theme_colors[i % len(theme_colors)] for i in range(len(data_arr))]
    set_point_colors(series, point_colors)

    # Data labels — value + category for donut
    set_chart_datalabels(chart._element, show_val=True, show_cat=True, show_pct=False)

def set_slide_background(slide, background):
    bg_type  = background.get('type','solid')
    bg_color = background.get('color','#1E3A5F')
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = safe_color(bg_color)

def generate(deck, out_path, image_paths=None):
    prs = Presentation()
    prs.slide_width  = I(SW)
    prs.slide_height = I(SH)

    blank_layout = prs.slide_layouts[6]

    for si, slide_data in enumerate(deck.get('slides', [])):
        slide = prs.slides.add_slide(blank_layout)

        bg = slide_data.get('background', {'type':'solid','color':'#1E3A5F'})
        set_slide_background(slide, bg)

        elements = slide_data.get('elements', [])
        elements_sorted = sorted(elements, key=lambda e: e.get('z', 2))
        for el in elements_sorted:
            try:
                render_element(slide, el, si)
            except Exception as err:
                print(f'[PPTX] Element error (slide {si}, {el.get("type")}): {err}', file=sys.stderr)

    prs.save(out_path)
    print(f'[PPTX] Saved {len(deck.get("slides",[]))} slides → {out_path}')

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: generate_pptx.py <deck_json> <output.pptx> [--images img1,img2]', file=sys.stderr)
        sys.exit(1)

    deck_arg = sys.argv[1]
    out_path = sys.argv[2]
    image_paths = []
    if '--images' in sys.argv:
        idx = sys.argv.index('--images')
        if idx + 1 < len(sys.argv):
            image_paths = [p for p in sys.argv[idx+1].split(',') if p]

    try:
        deck = json.loads(deck_arg)
    except json.JSONDecodeError:
        print(f'[PPTX] Invalid JSON input', file=sys.stderr)
        sys.exit(1)

    # Support legacy format (plain slides array)
    if isinstance(deck, list):
        deck = {'theme': {'name':'blue'}, 'slides': [{'id':f's{i}','background':{'type':'solid','color':'#1E3A5F'},'elements':[],'raw':s} for i,s in enumerate(deck)]}

    generate(deck, out_path, image_paths)
